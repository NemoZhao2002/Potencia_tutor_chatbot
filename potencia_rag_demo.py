import os
import logging
import gradio as gr
from pptx import Presentation
import fitz
from llama_index.llms.openai import OpenAI
from typing import Optional
from llama_index.core import (
    SimpleDirectoryReader,
    VectorStoreIndex,
    StorageContext,
    load_index_from_storage,
)
from llama_index.core.schema import Document, QueryBundle
from llama_index.core.tools import QueryEngineTool, BaseTool
from llama_index.core.query_engine import RouterQueryEngine, BaseQueryEngine
from llama_index.core.chat_engine.types import BaseChatEngine
from llama_index.core.selectors import PydanticSingleSelector
from llama_index.core.tools.types import AsyncBaseTool, ToolMetadata
from llama_index.core.base.response.schema import Response
from llama_index.core.callbacks.base import CallbackManager
# Logging setup
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

# Environment setup
os.environ['OPENAI_API_KEY'] = os.getenv('API_KEY')
VECTOR_STORAGE_DIR = "./vector"
DATA_STORAGE_DIR = "./data"
os.makedirs(DATA_STORAGE_DIR, exist_ok=True)

# Helper Functions
def get_file_names_in_directory(directory_path):
    """List all files in the directory"""
    file_paths = []
    names = []
    for dirpath, _, filenames in os.walk(directory_path):
        for filename in filenames:
            file_paths.append(os.path.join(dirpath, filename))
            names.append(filename)
    return file_paths, names

def load_pptx_text(file_path):
    """Extract text from a PPTX file."""
    pres = Presentation(file_path)
    pptx_text = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pptx_text.append(run.text)
    return pptx_text

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file using PyMuPDF."""
    text = []
    try:
        pdf_document = fitz.open(file_path)
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            text.append(page.get_text())
        pdf_document.close()
    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
    return "\n".join(text)

def load_any_documents(dirpath: str):
    """Load all documents from a directory."""
    documents = []
    for file_name in os.listdir(dirpath):
        file_path = os.path.join(dirpath, file_name)
        if file_name.endswith(".pptx"):
            pptx_text = load_pptx_text(file_path)
            documents.append(Document(text="\n".join(pptx_text), doc_id=file_name))
        elif file_name.endswith(".pdf"):
            pdf_text = extract_text_from_pdf(file_path)
            documents.append(Document(text=pdf_text, doc_id=file_name))
        else:
            documents += SimpleDirectoryReader(input_files=[file_path]).load_data()
    return documents

class ChatEngineAdapter(BaseQueryEngine):
    """
    Adapter to make ChatEngine compatible with QueryEngineTool.
    """

    def __init__(self, chat_engine: BaseChatEngine, callback_manager: CallbackManager = None):
        self.chat_engine = chat_engine
        self.callback_manager = callback_manager or CallbackManager([])

    def _query(self, query_bundle: QueryBundle) -> Response:
        """
        Implements the abstract `_query` method.
        Converts AgentChatResponse to Response.
        """
        # Extract query string
        query_str = query_bundle.query_str
        
        # Use the chat engine to process the query
        agent_response = self.chat_engine.chat(query_str)
        print(agent_response)
        # Convert AgentChatResponse to Response
        return Response(
            response=agent_response.response,  # Text content of the response
            metadata=agent_response.metadata,  # Any additional metadata
            source_nodes=agent_response.source_nodes  # If available
        )

    async def _aquery(self, query_bundle: QueryBundle) -> Response:
        """
        Implements the abstract `_aquery` method for asynchronous queries.
        Converts AgentChatResponse to Response.
        """
        # Extract query string
        query_str = query_bundle.query_str
        
        # Use the chat engine to process the query asynchronously
        agent_response = await self.chat_engine.achat(query_str)
        
        # Convert AgentChatResponse to Response
        return Response(
            response=agent_response.response,
            metadata=agent_response.metadata,
            source_nodes=agent_response.source_nodes
        )
        
    def _get_prompt_modules(self) -> dict:
        """
        Provide a default implementation for `_get_prompt_modules`.
        Returns an empty dictionary for compatibility.
        """
        return {}
    
class ChatEngineTool(QueryEngineTool):
    """
    A tool to wrap a ChatEngine using an adapter for compatibility with QueryEngineTool.
    """

    def __init__(self, chat_engine: BaseChatEngine, description: str, name: str = "ChatEngineTool"):
        # Wrap ChatEngine with an adapter to make it compatible
        query_engine_adapter = ChatEngineAdapter(chat_engine)

        # Create metadata for the tool
        metadata = ToolMetadata(
            name=name,
            description=description,
            return_direct=False  # Customize as needed
        )

        # Initialize the parent class with the adapter and metadata
        super().__init__(query_engine=query_engine_adapter, metadata=metadata)
        self.chat_engine = chat_engine

    @classmethod
    def from_defaults(cls, chat_engine: BaseChatEngine, description: str, name: str = "ChatEngineTool"):
        """
        Create a ChatEngineTool with default settings.
        """
        return cls(chat_engine=chat_engine, description=description, name=name)


# Load and Process Indices
def initialize_query_engines():
    index_list = []
    description_list = []
    query_engine_tools = []

    for dirpath, _, file_names in os.walk(DATA_STORAGE_DIR):
        if file_names:
            relative_path = os.path.relpath(dirpath, DATA_STORAGE_DIR)
            storage_dir = os.path.join(VECTOR_STORAGE_DIR, relative_path)
            os.makedirs(storage_dir, exist_ok=True)
            
            try:
                storage_context = StorageContext.from_defaults(persist_dir=storage_dir)
                index = load_index_from_storage(storage_context)
                description = storage_dir.split('/')[-1]
                index_list.append(index)
                description_list.append(description)
                logger.info(f"Loaded existing index for subdirectory {relative_path}.")
            except FileNotFoundError:
                logger.info(f"Index for {relative_path} not found. Creating a new index.")
                documents = load_any_documents(dirpath)
                index = VectorStoreIndex.from_documents(documents)
                index.storage_context.persist(storage_dir)
                index_list.append(index)
                description_list.append(relative_path)
            except Exception as e:
                logger.error(f"Error processing subdirectory {relative_path}: {e}")
    # return index_list[2].as_chat_engine()
    # Create query engine tools
    for description, index in zip(description_list, index_list):
        query_engine = index.as_query_engine(similarity_top_k=5)
        query_engine_tools.append(
            QueryEngineTool.from_defaults(
                query_engine=query_engine, 
                description=(f"Handles queries related to {description}.")
            )
        )
    
    return RouterQueryEngine(
        selector=PydanticSingleSelector.from_defaults(),
        query_engine_tools=query_engine_tools,
        verbose=True
    )

# # Initialize the RouterQueryEngine
router_query_engine = initialize_query_engines()
# # response = chat_engine.chat("Generate a lesson plan on verbs, use examples from previous include specific practice problems")
# # print(str(response))
# query = "Generate a lesson plan on verbs, use a lot of examples from previous include specific practice problems"
# response = router_chat_engine.query(query)
# print(f"Routerqueryengine: {str(response)}")
# print(f"context: ")
# for source_node in response.source_nodes[:1]:
#     # Get the content of the node
#     content = source_node.node.get_content()
#     # Get additional metadata if available
#     metadata = source_node.node.metadata
#     print(f"Content: {content}")
#     print(f"Metadata: {metadata}")
#     print("-" * 50)

# response = query_engine.query(query)
# print(f"single queryengine: {str(response)}")


# Gradio Function
def query(selected_option, custom_input):
    # Combine dropdown and text input into a single response
    input_query = custom_input.strip() if custom_input.strip() else selected_option
    response = router_query_engine.query(f"{input_query}, format the response")
    content = response.source_nodes[0].get_content()
    metadata = response.source_nodes[0].metadata
    
    return response, content, metadata

inputs = [
        gr.Dropdown(choices=["How much do the tutors get paid?", "What is the recommended structure for a tutoring session?", "Generate a lesson plan on irregular verbs for a beginner English class, include examples", "What forms of support does Potencia offer to tutors during the semester?"], label="Choose a question"),
        gr.Textbox(lines=1, placeholder="type your own question here", label="Custom Input")
    ]
# Gradio Interface
interface = gr.Interface(
    fn=query,
    inputs=inputs,
    outputs=[
        gr.Textbox(label="response"),
        gr.Textbox(label="Context"),
        gr.Textbox(label="Metadata"),
    ],
    title="RAG Query Interface",
    description="Ask questions about your documents and get accurate answers using the RAG framework."
)

# Launch the Gradio App
interface.launch()
